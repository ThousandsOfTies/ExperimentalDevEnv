#!/usr/bin/env python3
"""Generate intro presentation for ExperimentalDevEnv."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Color palette
C_BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
C_BG_MID    = RGBColor(0x1A, 0x2A, 0x4A)   # mid navy
C_ACCENT    = RGBColor(0x00, 0xC8, 0xFF)   # cyan
C_ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)   # purple
C_GREEN     = RGBColor(0x00, 0xE0, 0x9E)   # teal-green
C_ORANGE    = RGBColor(0xFF, 0x6B, 0x35)   # orange
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xB0, 0xC4, 0xDE)   # light blue-gray
C_YELLOW    = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def make_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── helper: filled rectangle ──────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        # set transparency via XML
        sp = shape._element
        solidFill = sp.find('.//' + pptx.oxml.ns.qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(pptx.oxml.ns.qn('a:srgbClr'))
            if srgbClr is not None:
                alpha_elem = etree.SubElement(srgbClr, pptx.oxml.ns.qn('a:alpha'))
                alpha_elem.set('val', str(alpha))
    return shape


# ── helper: text box ──────────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Hiragino Sans" if False else "Arial"
    return txBox


def add_text_para(slide, lines, left, top, width, height,
                  font_size=18, color=C_WHITE, align=PP_ALIGN.LEFT,
                  line_font_sizes=None, line_colors=None, line_bolds=None):
    """Multi-line text box with per-line control."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if not line:
            p.space_before = Pt(4)
            continue
        run = p.add_run()
        run.text = line
        fs = (line_font_sizes[i] if line_font_sizes and i < len(line_font_sizes) else font_size)
        cl = (line_colors[i]     if line_colors     and i < len(line_colors)     else color)
        bd = (line_bolds[i]      if line_bolds       and i < len(line_bolds)      else False)
        run.font.size = Pt(fs)
        run.font.color.rgb = cl
        run.font.bold = bd
        run.font.name = "Arial"
    return txBox


# ── slide background ──────────────────────────────────────────────────────────
def set_bg(slide, color=C_BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── accent bar (top) ──────────────────────────────────────────────────────────
def accent_bar(slide, color=C_ACCENT):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), color)


# =============================================================================
#  SLIDE 1: Title
# =============================================================================
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    # gradient left panel
    add_rect(slide, 0, 0, Inches(5.5), SLIDE_H, C_BG_MID)

    # accent lines
    add_rect(slide, 0, 0, Inches(0.15), SLIDE_H, C_ACCENT)
    add_rect(slide, Inches(5.35), 0, Inches(0.15), SLIDE_H, C_ACCENT2)

    # tag
    add_text(slide, "PoC / 技術提案", Inches(0.3), Inches(0.5),
             Inches(4.5), Inches(0.5), font_size=13, color=C_ACCENT)

    # main title
    add_text_para(slide,
        ["クラウドネイティブ", "組み込み開発環境"],
        Inches(0.3), Inches(1.2), Inches(4.8), Inches(2.0),
        font_size=36, color=C_WHITE, line_bolds=[True, True])

    # subtitle
    add_text(slide,
        "〜 AIと自動化で変わる\n  ハードウェア開発の新常識 〜",
        Inches(0.3), Inches(3.0), Inches(5.0), Inches(1.4),
        font_size=18, color=C_LIGHT, italic=True)

    # right side visual: pipeline boxes
    boxes = [
        (C_ACCENT,  "☁  GitHub Codespaces",  "クロスコンパイル環境"),
        (C_ACCENT2, "⚡  AWS EC2 Graviton",   "ARM64 クラウドシミュレーション"),
        (C_GREEN,   "🔧  Raspberry Pi 5",     "実機デプロイ（変更なし）"),
    ]
    for i, (col, title, sub) in enumerate(boxes):
        y = Inches(1.5 + i * 1.7)
        add_rect(slide, Inches(6.0), y, Inches(6.8), Inches(1.3), C_BG_MID)
        add_rect(slide, Inches(6.0), y, Inches(0.15), Inches(1.3), col)
        add_text(slide, title, Inches(6.3), y + Inches(0.15),
                 Inches(6.2), Inches(0.5), font_size=16, bold=True, color=col)
        add_text(slide, sub, Inches(6.3), y + Inches(0.65),
                 Inches(6.2), Inches(0.4), font_size=13, color=C_LIGHT)
        if i < 2:
            add_text(slide, "↓", Inches(9.2), y + Inches(1.3),
                     Inches(0.5), Inches(0.3), font_size=18, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)

    # date
    add_text(slide, "2026年5月", Inches(0.3), Inches(6.8),
             Inches(3), Inches(0.5), font_size=12, color=C_LIGHT)


# =============================================================================
#  SLIDE 2: 課題 (Problem)
# =============================================================================
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_ORANGE)

    add_text(slide, "01  課題", Inches(0.5), Inches(0.15),
             Inches(4), Inches(0.5), font_size=12, color=C_ORANGE, bold=True)
    add_text(slide, "従来の組み込み開発が\n抱えるボトルネック",
             Inches(0.5), Inches(0.5), Inches(7), Inches(1.2),
             font_size=30, bold=True, color=C_WHITE)

    problems = [
        (C_ORANGE,  "🔩  ハードウェア依存",
         "実機がなければ開発・テストできない\n"
         "→ 新人・リモートワーカーが参加できない"),
        (C_ACCENT2, "🔄  環境差異リスク",
         "開発者ごとに異なる OS / ツールチェーン\n"
         "→「自分の環境では動く」問題が頻発"),
        (C_YELLOW,  "🐢  CI/CD の未整備",
         "組み込みバイナリの自動テストが困難\n"
         "→ バグが実機フェーズまで発見されない"),
        (C_GREEN,   "💸  コスト・スケール",
         "テスト用ハードウェア調達・管理コストが高額\n"
         "→ 開発台数分のデバイスが必要"),
    ]

    for i, (col, title, body) in enumerate(problems):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.5 + col_idx * 6.3)
        y = Inches(2.0 + row_idx * 2.3)
        add_rect(slide, x, y, Inches(5.8), Inches(2.0), C_BG_MID)
        add_rect(slide, x, y, Inches(0.12), Inches(2.0), col)
        add_text(slide, title, x + Inches(0.25), y + Inches(0.1),
                 Inches(5.3), Inches(0.5), font_size=15, bold=True, color=col)
        add_text(slide, body, x + Inches(0.25), y + Inches(0.6),
                 Inches(5.3), Inches(1.2), font_size=13, color=C_LIGHT)


# =============================================================================
#  SLIDE 3: 業界トレンド (Industry Trends)
# =============================================================================
def slide_trends(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_ACCENT2)

    add_text(slide, "02  業界トレンド", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_ACCENT2, bold=True)
    add_text(slide, "AI・クラウドが変える\n組み込み開発の地平",
             Inches(0.5), Inches(0.5), Inches(8), Inches(1.2),
             font_size=30, bold=True, color=C_WHITE)

    trends = [
        (C_ACCENT,  "🤖  AI Coding Assistant の台頭",
         ["GitHub Copilot / Claude Code がドライバ・テストコードを自動生成",
          "AI はクラウド上で動作 → ハードウェアへの直接アクセス不可",
          "→ 仮想デバイス層が AI 活用の鍵になる"]),
        (C_ACCENT2, "☁  クラウド IDE の普及",
         ["GitHub Codespaces / Gitpod が標準化",
          "ブラウザだけで本格的な C/C++ クロスコンパイルが可能",
          "→ ローカル環境セットアップゼロ"]),
        (C_GREEN,   "🏭  デジタルツイン化の加速",
         ["AWS IoT TwinMaker / Azure Digital Twins が成長市場",
          "実機と同一ロジックをクラウドで先行検証するアプローチが主流に",
          "→ 本 PoC はその最小構成を実証"]),
        (C_YELLOW,  "📡  Edge AI × ARM の爆発的成長",
         ["Raspberry Pi 5 / NVIDIA Jetson / Apple M シリーズすべて ARM64",
          "クラウド（EC2 Graviton）と実機が同一 ISA",
          "→ 「一度ビルド、どこでも動く」が現実的に"]),
    ]

    for i, (col, title, bullets) in enumerate(trends):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.5 + col_idx * 6.3)
        y = Inches(2.0 + row_idx * 2.3)
        add_rect(slide, x, y, Inches(5.9), Inches(2.1), C_BG_MID)
        add_rect(slide, x, y, Inches(0.12), Inches(2.1), col)
        add_text(slide, title, x + Inches(0.25), y + Inches(0.08),
                 Inches(5.4), Inches(0.45), font_size=14, bold=True, color=col)
        bullet_text = "\n".join(f"• {b}" for b in bullets)
        add_text(slide, bullet_text, x + Inches(0.25), y + Inches(0.55),
                 Inches(5.4), Inches(1.4), font_size=11.5, color=C_LIGHT)


# =============================================================================
#  SLIDE 4: Solution Overview
# =============================================================================
def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_GREEN)

    add_text(slide, "03  ソリューション", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_GREEN, bold=True)
    add_text(slide, "ExperimentalDevEnv",
             Inches(0.5), Inches(0.5), Inches(10), Inches(0.7),
             font_size=32, bold=True, color=C_WHITE)
    add_text(slide,
        "「ハードウェアなしで始め、同じバイナリを実機に展開する」\nクラウドネイティブ組み込み開発 PoC",
        Inches(0.5), Inches(1.15), Inches(12), Inches(0.8),
        font_size=16, color=C_LIGHT, italic=True)

    # 3-column value props
    cols = [
        (C_ACCENT,  "Zero Setup",
         "GitHub Codespaces で\nクロスコンパイル環境を\n即時起動\n\nローカル環境構築不要"),
        (C_ACCENT2, "Cloud Simulate",
         "AWS EC2 Graviton で\nGPIO / I2C / SPI を仮想化\n\nブラウザで\nハードウェアを操作"),
        (C_GREEN,   "Real Deploy",
         "同一バイナリを\nRaspberry Pi 5 に展開\n\nコード変更ゼロで\n実機動作を確認"),
    ]
    for i, (col, title, body) in enumerate(cols):
        x = Inches(0.4 + i * 4.2)
        y = Inches(2.1)
        add_rect(slide, x, y, Inches(3.9), Inches(4.7), C_BG_MID)
        add_rect(slide, x, y, Inches(3.9), Inches(0.12), col)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.25),
                 Inches(3.5), Inches(0.6), font_size=20, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.2), y + Inches(0.9),
                 Inches(3.5), Inches(3.5), font_size=14, color=C_LIGHT,
                 align=PP_ALIGN.CENTER)

    # arrow connector between boxes
    for i in range(2):
        x = Inches(4.25 + i * 4.2)
        add_text(slide, "→", x, Inches(4.1), Inches(0.5), Inches(0.6),
                 font_size=24, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_text(slide,
        "★  開発から実機検証まで全工程をクラウドで完結  ★",
        Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
        font_size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


# =============================================================================
#  SLIDE 5: Architecture
# =============================================================================
def slide_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_ACCENT)

    add_text(slide, "04  アーキテクチャ", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_ACCENT, bold=True)
    add_text(slide, "3層パイプライン構成",
             Inches(0.5), Inches(0.5), Inches(8), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)

    layers = [
        (C_ACCENT,  "Tier 1  ビルド層",
         "GitHub Codespaces  (x86_64, Ubuntu 24.04)",
         ["aarch64-linux-gnu-gcc でクロスコンパイル",
          "libfuse3, linux-headers 含む完全な組み込み toolchain",
          "VSCode タスクで one-click ビルド & デプロイ"]),
        (C_ACCENT2, "Tier 2  シミュレーション層",
         "AWS EC2 Graviton  (ARM64, Ubuntu 24.04)",
         ["LD_PRELOAD shim で GPIO / SPI をインターセプト",
          "CUSE（Userspace キャラクタデバイス）で I2C をエミュレート",
          "Python WebSocket ブリッジ → ブラウザのインタラクティブ HW パネル"]),
        (C_GREEN,   "Tier 3  実機層",
         "Raspberry Pi 5  (ARM64, Raspberry Pi OS)",
         ["同一バイナリを scp でデプロイ（変更なし）",
          "lgpio で実 GPIO / I2C / SPI デバイスを直接制御",
          "VL53L0X センサー・MFRC-522 RFID・SSD1306 OLED 動作確認済"]),
    ]

    for i, (col, tier, platform, bullets) in enumerate(layers):
        y = Inches(1.45 + i * 1.9)
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(1.7), C_BG_MID)
        add_rect(slide, Inches(0.4), y, Inches(0.12), Inches(1.7), col)
        add_text(slide, tier, Inches(0.65), y + Inches(0.08),
                 Inches(3.5), Inches(0.45), font_size=14, bold=True, color=col)
        add_text(slide, platform, Inches(4.0), y + Inches(0.08),
                 Inches(8.6), Inches(0.45), font_size=13, color=C_LIGHT)
        bullet_text = "   ".join(f"• {b}" for b in bullets)
        add_text(slide, bullet_text, Inches(0.65), y + Inches(0.6),
                 Inches(12.1), Inches(0.9), font_size=11.5, color=C_LIGHT)
        if i < 2:
            add_text(slide, "↓  cross-compile & scp", Inches(5.5), y + Inches(1.7),
                     Inches(4), Inches(0.25), font_size=11, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)


# =============================================================================
#  SLIDE 6: Key Technologies
# =============================================================================
def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_ACCENT2)

    add_text(slide, "05  技術詳細", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_ACCENT2, bold=True)
    add_text(slide, "ハードウェア仮想化の仕組み",
             Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)

    techs = [
        (C_ORANGE,  "GPIO シム\ngpio_shim.so",
         "LD_PRELOAD フック\n/dev/gpiochip0 ioctl を\nインターセプト\n\nUnix ソケット経由で\nLED 状態を送信\nボタン入力を注入"),
        (C_ACCENT2, "I2C CUSE スタブ\ncuse_i2c",
         "FUSE3 で\n/dev/i2c-1 仮想デバイス生成\n\nアドレスで振り分け:\n0x29 → VL53L0X\n0x3C → SSD1306"),
        (C_ACCENT,  "SPI シム\nspi_shim.so",
         "LD_PRELOAD フック\n/dev/spidev0.0 ioctl\n\nMFRC-522 レジスタ\nプロトコルを\n完全エミュレート"),
        (C_GREEN,   "Web ブリッジ\nbridge.py",
         "Python asyncio\nWebSocket ↔\nUnix ソケット\n\nブラウザで\nリアルタイム\nHW 可視化"),
    ]

    for i, (col, title, body) in enumerate(techs):
        x = Inches(0.4 + i * 3.1)
        y = Inches(1.5)
        add_rect(slide, x, y, Inches(2.85), Inches(5.5), C_BG_MID)
        add_rect(slide, x, y, Inches(2.85), Inches(0.1), col)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.2),
                 Inches(2.55), Inches(0.9), font_size=13, bold=True, color=col)
        add_text(slide, body, x + Inches(0.15), y + Inches(1.1),
                 Inches(2.55), Inches(4.1), font_size=12, color=C_LIGHT)

    # bottom note
    add_text(slide,
        "アプリケーションコードは EC2 シム / RasPi5 実機で完全に同一。LD_PRELOAD や CUSE は透過的に置き換わる。",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
        font_size=12, color=C_GREEN, align=PP_ALIGN.CENTER)


# =============================================================================
#  SLIDE 7: Results / Validation
# =============================================================================
def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_GREEN)

    add_text(slide, "06  実績・検証結果", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_GREEN, bold=True)
    add_text(slide, "PoC で実証済みの機能",
             Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)

    # Validation table
    headers = ["コンポーネント", "インターフェース", "EC2 シム", "RasPi5 実機", "ステータス"]
    rows = [
        ["LED + ボタン", "GPIO",  "✅ 完全シミュレーション", "✅ 実機テスト済",    "完了"],
        ["VL53L0X センサー", "I2C", "✅ レジスタエミュレーション", "✅ 実デバイス動作",  "完了"],
        ["MFRC-522 RFID", "SPI",  "✅ レジスタエミュレーション", "✅ カード読み取り",  "完了"],
        ["SSD1306 OLED", "I2C",  "✅ フレームバッファ→Canvas", "⏳ モジュール待ち", "進行中"],
        ["Web パネル", "WebSocket", "✅ リアルタイム操作",       "—",              "完了"],
    ]

    col_widths = [Inches(2.5), Inches(1.8), Inches(3.0), Inches(3.0), Inches(1.5)]
    col_x = [Inches(0.4), Inches(2.9), Inches(4.7), Inches(7.7), Inches(10.7)]
    row_h = Inches(0.7)
    header_y = Inches(1.5)

    # header row
    add_rect(slide, Inches(0.4), header_y, Inches(12.4), row_h, C_BG_MID)
    add_rect(slide, Inches(0.4), header_y, Inches(12.4), Inches(0.06), C_GREEN)
    for j, (h, w, x) in enumerate(zip(headers, col_widths, col_x)):
        add_text(slide, h, x + Inches(0.05), header_y + Inches(0.15),
                 w, Inches(0.4), font_size=12, bold=True, color=C_GREEN)

    for i, row in enumerate(rows):
        y = header_y + row_h * (i + 1)
        bg = C_BG_MID if i % 2 == 0 else RGBColor(0x14, 0x20, 0x38)
        add_rect(slide, Inches(0.4), y, Inches(12.4), row_h, bg)
        for j, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            col = C_GREEN if cell.startswith("✅") else (C_YELLOW if cell.startswith("⏳") else C_WHITE)
            if j == 4:
                col = C_GREEN if cell == "完了" else C_YELLOW
            add_text(slide, cell, x + Inches(0.05), y + Inches(0.18),
                     w, Inches(0.4), font_size=11.5, color=col)

    # key stat callouts
    stats = [
        (C_ACCENT,  "~1,700 行", "C コード"),
        (C_ACCENT2, "3 種",     "HW I/F\n(GPIO/I2C/SPI)"),
        (C_GREEN,   "0 行",     "実機向け\nコード変更"),
    ]
    for i, (col, num, label) in enumerate(stats):
        x = Inches(0.4 + i * 4.2)
        y = Inches(6.45)
        add_rect(slide, x, y, Inches(3.9), Inches(0.9), C_BG_MID)
        add_rect(slide, x, y, Inches(0.1), Inches(0.9), col)
        add_text(slide, num, x + Inches(0.2), y + Inches(0.05),
                 Inches(1.5), Inches(0.5), font_size=22, bold=True, color=col)
        add_text(slide, label, x + Inches(1.7), y + Inches(0.1),
                 Inches(2.0), Inches(0.7), font_size=11, color=C_LIGHT)


# =============================================================================
#  SLIDE 8: AI Integration Opportunity
# =============================================================================
def slide_ai(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_YELLOW)

    add_text(slide, "07  AI 活用の可能性", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_YELLOW, bold=True)
    add_text(slide, "Claude Code × 組み込み開発",
             Inches(0.5), Inches(0.5), Inches(10), Inches(0.7),
             font_size=30, bold=True, color=C_WHITE)
    add_text(slide,
        "本環境はクラウド上の AI との親和性が高く、AI-Driven 開発フローへの拡張が自然に行える",
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
        font_size=15, color=C_LIGHT, italic=True)

    opportunities = [
        (C_ACCENT,  "💡  ドライバコード生成",
         "Claude Code が I2C / SPI デバイスのドライバを\n自動生成 → CUSE スタブで即テスト\n実機なしで動作確認が完結"),
        (C_ACCENT2, "🧪  自動テストシナリオ",
         "AI がシミュレーション環境に仮想センサー値を注入\n境界値・異常値を自動生成してテスト\nCI/CD パイプラインに組み込み可能"),
        (C_GREEN,   "📊  ハードウェアトレース解析",
         "GPIO / I2C / SPI のトレースログを AI に渡し\nバグ原因を自動推定・修正提案\nデバッグ工数を大幅削減"),
        (C_ORANGE,  "🚀  コード品質ゲート",
         "PR 時に AI がドライバコードをレビュー\nメモリ安全性・タイミング問題を静的解析\nハードウェア起因のバグを事前に防止"),
    ]

    for i, (col, title, body) in enumerate(opportunities):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.4 + col_idx * 6.4)
        y = Inches(2.0 + row_idx * 2.4)
        add_rect(slide, x, y, Inches(6.0), Inches(2.1), C_BG_MID)
        add_rect(slide, x, y, Inches(0.12), Inches(2.1), col)
        add_text(slide, title, x + Inches(0.25), y + Inches(0.1),
                 Inches(5.5), Inches(0.5), font_size=15, bold=True, color=col)
        add_text(slide, body, x + Inches(0.25), y + Inches(0.65),
                 Inches(5.5), Inches(1.3), font_size=12.5, color=C_LIGHT)

    add_text(slide,
        "→ AI は「クラウドに閉じた存在」だからこそ、クラウドで完結する本仮想 HW 環境と相性が最高",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
        font_size=13, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)


# =============================================================================
#  SLIDE 9: Roadmap / Next Steps
# =============================================================================
def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent_bar(slide, C_ACCENT)

    add_text(slide, "08  ロードマップ", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_ACCENT, bold=True)
    add_text(slide, "今後の展開・拡張計画",
             Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)

    phases = [
        ("Phase 1", "Now  ✅ PoC 完了", C_GREEN, [
            "GPIO / I2C / SPI シミュレーション確立",
            "EC2 Graviton 上でブラウザ HW パネル稼働",
            "RasPi5 実機デプロイ（同一バイナリ）検証済",
        ]),
        ("Phase 2", "Next  🔧 機能拡張", C_ACCENT, [
            "SSD1306 OLED シミュレーション完成",
            "LCD HAT (ST7789) シム追加",
            "CI/CD パイプライン統合（GitHub Actions）",
        ]),
        ("Phase 3", "Future  🚀 AI 統合", C_ACCENT2, [
            "Claude Code によるドライバ自動生成フロー",
            "AI 自動テストシナリオ生成",
            "複数デバイス対応（Jetson Nano, BeagleBone）",
        ]),
        ("Phase 4", "Vision  🌐 製品化", C_YELLOW, [
            "SaaS 型仮想 HW 環境プラットフォーム",
            "企業向け組み込み CI/CD ソリューション",
            "エッジ AI デバイス開発の標準環境へ",
        ]),
    ]

    for i, (phase, title, col, items) in enumerate(phases):
        x = Inches(0.4 + i * 3.1)
        y = Inches(1.5)
        add_rect(slide, x, y, Inches(2.9), Inches(5.6), C_BG_MID)
        add_rect(slide, x, y, Inches(2.9), Inches(0.1), col)

        add_text(slide, phase, x + Inches(0.15), y + Inches(0.15),
                 Inches(2.6), Inches(0.4), font_size=11, color=col, bold=True)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.5),
                 Inches(2.6), Inches(0.6), font_size=12.5, bold=True, color=C_WHITE)

        for j, item in enumerate(items):
            add_text(slide, f"• {item}", x + Inches(0.15),
                     y + Inches(1.25 + j * 1.2),
                     Inches(2.6), Inches(1.1), font_size=11.5, color=C_LIGHT)

    # timeline arrow
    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.4), Inches(0.12), C_BG_MID)
    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.4), Inches(0.12), C_ACCENT)
    add_text(slide, "→ 時間軸", Inches(11.5), Inches(6.95),
             Inches(1.8), Inches(0.35), font_size=12, color=C_ACCENT)


# =============================================================================
#  SLIDE 10: Summary / CTA
# =============================================================================
def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # full-width gradient top
    add_rect(slide, 0, 0, SLIDE_W, Inches(2.8), C_BG_MID)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), C_ACCENT)
    add_rect(slide, 0, Inches(2.8), SLIDE_W, Inches(0.06), C_ACCENT2)

    add_text(slide, "09  まとめ", Inches(0.5), Inches(0.15),
             Inches(5), Inches(0.5), font_size=12, color=C_ACCENT, bold=True)
    add_text(slide, "ExperimentalDevEnv が提供する価値",
             Inches(0.5), Inches(0.55), Inches(12), Inches(0.8),
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide,
        "クラウドで完結する組み込み開発パイプラインにより、\nハードウェアへの依存を排除しながら本物の組み込みソフトウェアを開発・検証できる",
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.1),
        font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

    values = [
        (C_ACCENT,  "🕐  開発速度 10×",   "ハード依存ゼロで\n即日開発開始"),
        (C_ACCENT2, "🔒  品質向上",       "クラウド CI/CD で\n自動テスト"),
        (C_GREEN,   "💰  コスト削減",     "テスト機材費・\n管理コストを圧縮"),
        (C_ORANGE,  "🤖  AI 対応",        "Claude Code と\nシームレスに統合"),
    ]

    for i, (col, title, body) in enumerate(values):
        x = Inches(0.5 + i * 3.1)
        y = Inches(3.1)
        add_rect(slide, x, y, Inches(2.85), Inches(2.8), C_BG_MID)
        add_rect(slide, x, y, Inches(2.85), Inches(0.1), col)
        add_text(slide, title, x + Inches(0.15), y + Inches(0.2),
                 Inches(2.55), Inches(0.6), font_size=16, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        add_text(slide, body, x + Inches(0.15), y + Inches(0.9),
                 Inches(2.55), Inches(1.5), font_size=13, color=C_LIGHT,
                 align=PP_ALIGN.CENTER)

    add_text(slide,
        "本 PoC をベースに、組み込み開発の新しいスタンダードを共に構築しましょう",
        Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6),
        font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide,
        "ExperimentalDevEnv  |  github.com/thousandsofties/experimentaldevenv",
        Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
        font_size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)


# =============================================================================
#  MAIN
# =============================================================================
def main():
    prs = make_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_trends(prs)
    slide_solution(prs)
    slide_arch(prs)
    slide_tech(prs)
    slide_results(prs)
    slide_ai(prs)
    slide_roadmap(prs)
    slide_summary(prs)

    out = "/home/user/ExperimentalDevEnv/docs/ExperimentalDevEnv_Intro.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
